from flask import Flask, request, jsonify
from flask_cors import CORS
import os
from pptx import Presentation

app = Flask(__name__)

# CORS設定（全てのオリジンを許可）
CORS(app, resources={r"/upload": {"origins": "*"}})

UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def rewrite_text(text):
    # 非AIによる簡易な丁寧語変換（例）
    replacements = {
        "です": "でございます",
        "ます": "いたします",
        "して": "いたしまして",
        "ありがとう": "ありがとうございます",
        "ください": "お願いいたします",
        "ない": "ございません",
        "できる": "可能です",
        "いる": "おります",
        "わかる": "理解いたします"
    }
    for word, polite in replacements.items():
        text = text.replace(word, polite)
    return text

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if not file.filename.endswith('.pptx'):
        return jsonify({'error': 'Invalid file format. Only .pptx files are allowed'}), 400

    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(file_path)

    try:
        ppt_text = extract_text_from_ppt(file_path)
        refined_text = rewrite_text(ppt_text)
        return jsonify({'ppt_text': refined_text})
    except Exception as e:
        return jsonify({'error': f'Processing failed: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5001, host='0.0.0.0', use_reloader=False)
